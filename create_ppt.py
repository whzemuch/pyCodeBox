import argparse
from pptx import Presentation
from pptx.util import Inches
from glob import glob

def slide_add_pic(prs:Presentation, pic:str)-> Presentation:
    """Add picture to a slide in a given presentation

    Args:
        prs (Presentation): a Presentation object
        pic (str): the file path of the picture

    Returns:
        Presentation: an updated Presentation object
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left, top, width = Inches(0.3), Inches(0.3), Inches(9.5)
    slide.shapes.add_picture(pic, left, top, width)

    return prs

def main():
    parser = argparse.ArgumentParser(description='Create a ppt file from images in a folder.')
    parser.add_argument('folder_path', type=str, help='The path to the folder containing the images')
    parser.add_argument('image_ext', type=str, help='The extension of the image file, for example: png, jpg')
    parser.add_argument('ppt_name', type=str, help='The file name for the ppt file')

    args = parser.parse_args()

    prs = Presentation()

    # get the path for all photos
    photo_files = glob(args.folder_path + '/*.' + args.image_ext)

    # insert each photo onto its own slide
    [slide_add_pic(prs, photo) for photo in sorted(photo_files)]

    # save the presentation
    prs.save(args.ppt_name)

if (__name__ == "__main__"):
    main()
