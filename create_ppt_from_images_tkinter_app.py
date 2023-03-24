import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
from glob import glob

class MainWindow(tk.Frame):
    def __init__(self, master=None):
        super().__init__(master)
        self.master = master
        self.master.title("Create PPT from Images")
        self.pack()

        self.folder_button = tk.Button(self, text="Select Folder", command=self.select_folder)
        self.folder_button.pack()

        self.folder_path = tk.StringVar()
        self.folder_path_label = tk.Label(self, textvariable=self.folder_path)
        self.folder_path_label.pack()

        self.image_ext_label = tk.Label(self, text="Image Extension:")
        self.image_ext_label.pack()

        self.image_ext = tk.StringVar()
        self.image_ext_entry = tk.Entry(self, textvariable=self.image_ext)
        self.image_ext_entry.pack()

        self.ppt_name_label = tk.Label(self, text="PPT Name:")
        self.ppt_name_label.pack()

        self.ppt_name = tk.StringVar()
        self.ppt_name_entry = tk.Entry(self, textvariable=self.ppt_name)
        self.ppt_name_entry.pack()

        self.save_button = tk.Button(self, text="Save PPT", command=self.save_ppt)
        self.save_button.pack()

    def select_folder(self):
        folder_path = filedialog.askdirectory()
        if folder_path:
            self.folder_path.set(folder_path)

    def save_ppt(self):
        folder_path = self.folder_path.get()
        image_ext = self.image_ext.get()
        ppt_name = self.ppt_name.get()

        if not folder_path or not image_ext or not ppt_name:
            messagebox.showerror("Error", "Please fill in all fields.")
            return

        image_files = glob(os.path.join(folder_path, f"*.{image_ext}"))

        if not image_files:
            messagebox.showerror("Error", f"No images with extension '{image_ext}' found in the selected folder.")
            return

        prs = Presentation()

        for image_file in sorted(image_files):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            left, top, width = Inches(0.3), Inches(0.3), Inches(9.5)
            slide.shapes.add_picture(image_file, left, top, width)

        ppt_path = os.path.join(folder_path, ppt_name)
        prs.save(ppt_path)

        messagebox.showinfo("Success", f"PPT '{ppt_name}' created successfully in '{folder_path}'.")

if __name__ == "__main__":
    root = tk.Tk()
    app = MainWindow(master=root)
    app.mainloop()
